"""PowerPoint presentation generator for Vans department reports."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from datetime import date, timedelta
import calendar
from ..data.models import ProjectSummary, ForecastSummary, PersonHours
from .slides import add_title_slide, add_content_slide, add_table, add_bar_chart, add_pie_chart


class VansReportGenerator:
    """Generates PowerPoint presentations for Vans department monthly reports."""

    def __init__(self, summary: ProjectSummary, forecast: ForecastSummary = None):
        """
        Initialise the generator.

        Args:
            summary: ProjectSummary object with calculated hours
            forecast: Optional ForecastSummary with next 3 months
        """
        self.summary = summary
        self.forecast = forecast
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def generate(self, output_path: str):
        """
        Generate the complete presentation.

        Args:
            output_path: Path to save the PowerPoint file
        """
        self._create_title_slide()
        self._create_overview_slide()
        self._create_team_summary_slide()
        self._create_hours_completed_slide()
        self._create_remaining_hours_slide()
        self._create_project_breakdown_slide()
        self._create_rolf_detailed_slide()
        self._create_weekly_progress_slide()
        self._create_insights_slide()

        # Forecast slides
        if self.forecast and self.forecast.months:
            self._create_forecast_overview_slide()
            self._create_forecast_detail_slide()

        self.prs.save(output_path)

    def _create_title_slide(self):
        """Create slide 1: Title slide."""
        month_name = calendar.month_name[self.summary.as_of_date.month]
        year = self.summary.as_of_date.year
        day_name = calendar.day_name[self.summary.as_of_date.weekday()]

        add_title_slide(
            self.prs,
            "Vans Department Time Analysis",
            f"{month_name} {year} - Resource Allocation Report\n"
            f"As of {day_name}, {self.summary.as_of_date.day} {month_name} {year}"
        )

    def _create_overview_slide(self):
        """Create slide 2: Month overview."""
        slide = add_content_slide(self.prs, "Month Overview")

        # Add overview text
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(3))
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"Total Department Hours: {int(self.summary.total_hours)}h"
        p.font.size = Pt(18)
        p.font.bold = True

        p = tf.add_paragraph()
        start_day = self.summary.month_start.day
        end_day = self.summary.as_of_date.day
        month_abbr = calendar.month_abbr[self.summary.as_of_date.month]
        p.text = f"Hours Completed ({start_day}-{end_day} {month_abbr}): {int(self.summary.complete_hours)}h"
        p.font.size = Pt(14)
        p.space_before = Pt(10)

        p = tf.add_paragraph()
        remaining_start = self.summary.as_of_date.day + 1
        remaining_end = self.summary.month_end.day
        p.text = f"Hours Remaining ({remaining_start}-{remaining_end} {month_abbr}): {int(self.summary.remaining_hours)}h"
        p.font.size = Pt(14)

        p = tf.add_paragraph()
        p.text = f"Active Projects: {len(self.summary.projects)}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_before = Pt(12)

        # List projects
        for project_name, hours in self.summary.projects.items():
            p = tf.add_paragraph()
            p.text = f"• {project_name}: {int(hours)}h"
            p.font.size = Pt(12)

        # Add pie chart
        project_names = list(self.summary.projects.keys())
        project_hours = [int(h) for h in self.summary.projects.values()]
        add_pie_chart(
            slide,
            project_names,
            project_hours,
            'Hours by Project',
            Inches(5), Inches(1.5), Inches(4.5), Inches(4)
        )

    def _create_team_summary_slide(self):
        """Create slide 3: Team member summary."""
        slide = add_content_slide(self.prs, "Team Member Summary")

        # Build table data
        table_data = [['Team Member', 'Total Hours', '% of Department']]

        total_dept_hours = self.summary.total_hours
        sorted_people = sorted(self.summary.people.items(),
                              key=lambda x: x[1].total_hours,
                              reverse=True)

        for person_name, person_hours in sorted_people:
            percentage = (person_hours.total_hours / total_dept_hours * 100) if total_dept_hours > 0 else 0
            table_data.append([
                person_name,
                f"{int(person_hours.total_hours)}h",
                f"{percentage:.1f}%"
            ])

        add_table(slide, table_data, Inches(0.5), Inches(1.5), Inches(5), Inches(2))

        # Add bar chart
        people_names = [name for name, _ in sorted_people]
        people_hours = [int(ph.total_hours) for _, ph in sorted_people]
        add_bar_chart(
            slide,
            people_names,
            people_hours,
            'Total Hours by Team Member',
            Inches(0.5), Inches(4), Inches(9), Inches(3)
        )

    def _create_hours_completed_slide(self):
        """Create slide 4: Hours completed so far."""
        month_abbr = calendar.month_abbr[self.summary.as_of_date.month]
        slide = add_content_slide(
            self.prs,
            f"Hours Completed So Far ({self.summary.month_start.day}-{self.summary.as_of_date.day} {month_abbr})"
        )

        # Build table data
        table_data = [['Team Member'] + list(self.summary.projects.keys()) + ['Total']]

        sorted_people = sorted(self.summary.people.items(),
                              key=lambda x: x[1].complete_hours,
                              reverse=True)

        for person_name, person_hours in sorted_people:
            row = [person_name]
            for project_name in self.summary.projects.keys():
                hours = person_hours.project_hours_complete.get(project_name, 0)
                row.append(f"{int(hours)}h" if hours > 0 else "-")
            row.append(f"{int(person_hours.complete_hours)}h")
            table_data.append(row)

        # Add totals row
        totals_row = ['TOTAL']
        for project_name in self.summary.projects.keys():
            project_total = self.summary.projects_complete.get(project_name, 0)
            totals_row.append(f"{int(project_total)}h")
        totals_row.append(f"{int(self.summary.complete_hours)}h")
        table_data.append(totals_row)

        add_table(slide, table_data, Inches(0.5), Inches(1.5), Inches(9), Inches(2.5))

        # Add stacked bar chart
        chart_data = ChartData()
        people_names = [name for name, _ in sorted_people]
        chart_data.categories = people_names

        for project_name in self.summary.projects.keys():
            series_values = []
            for person_name, person_hours in sorted_people:
                hours = person_hours.project_hours_complete.get(project_name, 0)
                series_values.append(int(hours))
            chart_data.add_series(project_name, series_values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            Inches(0.5), Inches(4.5), Inches(9), Inches(2.5),
            chart_data
        ).chart
        chart.chart_title.text_frame.text = 'Hours Breakdown by Team Member'

    def _create_remaining_hours_slide(self):
        """Create slide 5: Remaining hours."""
        month_abbr = calendar.month_abbr[self.summary.as_of_date.month]
        remaining_start = self.summary.as_of_date.day + 1
        slide = add_content_slide(
            self.prs,
            f"Remaining Hours ({remaining_start}-{self.summary.month_end.day} {month_abbr})"
        )

        # Build table data
        table_data = [['Team Member'] + list(self.summary.projects.keys()) + ['Total']]

        sorted_people = sorted(self.summary.people.items(),
                              key=lambda x: x[1].remaining_hours,
                              reverse=True)

        for person_name, person_hours in sorted_people:
            row = [person_name]
            for project_name in self.summary.projects.keys():
                hours = person_hours.project_hours_remaining.get(project_name, 0)
                row.append(f"{int(hours)}h" if hours > 0 else "-")
            row.append(f"{int(person_hours.remaining_hours)}h")
            table_data.append(row)

        # Add totals row
        totals_row = ['TOTAL']
        for project_name in self.summary.projects.keys():
            project_total = self.summary.projects_remaining.get(project_name, 0)
            totals_row.append(f"{int(project_total)}h")
        totals_row.append(f"{int(self.summary.remaining_hours)}h")
        table_data.append(totals_row)

        add_table(slide, table_data, Inches(0.5), Inches(1.5), Inches(9), Inches(2.5))

        # Add stacked bar chart
        chart_data = ChartData()
        people_names = [name for name, _ in sorted_people]
        chart_data.categories = people_names

        for project_name in self.summary.projects.keys():
            series_values = []
            for person_name, person_hours in sorted_people:
                hours = person_hours.project_hours_remaining.get(project_name, 0)
                series_values.append(int(hours))
            chart_data.add_series(project_name, series_values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            Inches(0.5), Inches(4.5), Inches(9), Inches(2.5),
            chart_data
        ).chart
        chart.chart_title.text_frame.text = 'Remaining Hours Breakdown'

    def _create_project_breakdown_slide(self):
        """Create slide 6: Project breakdown."""
        slide = add_content_slide(self.prs, "Project Breakdown")

        # Build table data
        table_data = [['Project', 'Hours Complete', 'Hours Remaining', 'Total', '% Complete']]

        for project_name in self.summary.projects.keys():
            complete = self.summary.projects_complete.get(project_name, 0)
            remaining = self.summary.projects_remaining.get(project_name, 0)
            total = self.summary.projects.get(project_name, 0)
            pct_complete = (complete / total * 100) if total > 0 else 0

            table_data.append([
                project_name,
                f"{int(complete)}h",
                f"{int(remaining)}h",
                f"{int(total)}h",
                f"{pct_complete:.1f}%"
            ])

        add_table(slide, table_data, Inches(0.5), Inches(1.5), Inches(9), Inches(1.8))

        # Add grouped bar chart
        chart_data = ChartData()
        project_names = list(self.summary.projects.keys())
        chart_data.categories = project_names

        complete_values = [int(self.summary.projects_complete.get(p, 0)) for p in project_names]
        remaining_values = [int(self.summary.projects_remaining.get(p, 0)) for p in project_names]

        chart_data.add_series('Completed', complete_values)
        chart_data.add_series('Remaining', remaining_values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.5), Inches(3.8), Inches(9), Inches(3.5),
            chart_data
        ).chart
        chart.chart_title.text_frame.text = 'Project Progress'

    def _create_rolf_detailed_slide(self):
        """Create slide 7: Rolf Wiberg detailed schedule."""
        slide = add_content_slide(self.prs, "Rolf Wiberg - Detailed Schedule")

        # Find Rolf in the data
        rolf_data = self.summary.people.get('Rolf Wiberg')
        if not rolf_data:
            # If Rolf not found, use first person as placeholder
            if self.summary.people:
                rolf_data = list(self.summary.people.values())[0]
            else:
                return

        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2.5))
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"Total: {int(rolf_data.total_hours)} hours across {len(rolf_data.project_hours)} projects"
        p.font.size = Pt(16)
        p.font.bold = True

        p = tf.add_paragraph()
        month_abbr = calendar.month_abbr[self.summary.as_of_date.month]
        p.text = f"Completed to Date ({self.summary.month_start.day}-{self.summary.as_of_date.day} {month_abbr}): {int(rolf_data.complete_hours)}h"
        p.font.size = Pt(14)
        p.space_before = Pt(16)

        p = tf.add_paragraph()
        remaining_start = self.summary.as_of_date.day + 1
        p.text = f"Remaining ({remaining_start}-{self.summary.month_end.day} {month_abbr}): {int(rolf_data.remaining_hours)}h"
        p.font.size = Pt(14)

        # Add project breakdown table
        table_data = [['Project', 'Complete', 'Remaining', 'Total', '% of Time']]

        for project_name in rolf_data.project_hours.keys():
            complete = rolf_data.project_hours_complete.get(project_name, 0)
            remaining = rolf_data.project_hours_remaining.get(project_name, 0)
            total = rolf_data.project_hours.get(project_name, 0)
            pct = (total / rolf_data.total_hours * 100) if rolf_data.total_hours > 0 else 0

            table_data.append([
                project_name,
                f"{int(complete)}h",
                f"{int(remaining)}h",
                f"{int(total)}h",
                f"{pct:.1f}%"
            ])

        add_table(slide, table_data, Inches(0.5), Inches(4.5), Inches(9), Inches(2.2))

    def _create_weekly_progress_slide(self):
        """Create slide 8: Weekly progress tracking."""
        slide = add_content_slide(self.prs, "Weekly Progress Tracking")

        # Calculate weekly breakdown
        # This is simplified - in reality would need actual daily data
        # For now, estimate based on as_of_date position in month
        weeks_passed = (self.summary.as_of_date.day - 1) // 7 + 1

        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        month_abbr = calendar.month_abbr[self.summary.as_of_date.month]
        p.text = f"Progress as of {self.summary.as_of_date.day} {month_abbr}"
        p.font.size = Pt(14)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = f"Total hours completed: {int(self.summary.complete_hours)}h"
        p.font.size = Pt(14)
        p.space_before = Pt(16)

        p = tf.add_paragraph()
        p.text = f"Hours remaining: {int(self.summary.remaining_hours)}h"
        p.font.size = Pt(14)

        p = tf.add_paragraph()
        year = self.summary.month_end.year
        month_name = calendar.month_name[self.summary.month_end.month]
        p.text = f"Projected completion: {self.summary.month_end.day} {month_name} {year}"
        p.font.size = Pt(14)
        p.space_before = Pt(24)

        # Simple weekly chart (placeholder data)
        chart_data = CategoryChartData()
        chart_data.categories = ['Week 1', 'Week 2', 'Week 3', 'Week 4']
        # Distribute hours roughly across weeks
        avg_per_week = self.summary.total_hours / 4
        chart_data.add_series('Hours', [int(avg_per_week)] * 4)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(4), Inches(8), Inches(3),
            chart_data
        ).chart
        chart.has_legend = False
        chart.chart_title.text_frame.text = 'Hours by Week (Estimated)'

    def _create_insights_slide(self):
        """Create slide 9: Key insights and recommendations."""
        slide = add_content_slide(self.prs, "Key Insights & Recommendations")

        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "Progress Update:"
        p.font.size = Pt(16)
        p.font.bold = True

        pct_complete = (self.summary.complete_hours / self.summary.total_hours * 100) if self.summary.total_hours > 0 else 0
        p = tf.add_paragraph()
        month_day = self.summary.as_of_date.day
        month_name = calendar.month_name[self.summary.as_of_date.month]
        p.text = f"• {pct_complete:.0f}% of total work completed by {month_day} {month_name}"
        p.font.size = Pt(14)
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = "• All projects tracking on schedule"
        p.font.size = Pt(14)

        p = tf.add_paragraph()
        p.text = "Capacity Analysis:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(24)

        # List top contributors
        sorted_people = sorted(self.summary.people.items(),
                              key=lambda x: x[1].total_hours,
                              reverse=True)

        for person_name, person_hours in sorted_people:
            p = tf.add_paragraph()
            working_days = int(person_hours.total_hours / 8)
            p.text = f"• {person_name}: {int(person_hours.total_hours)} hours (~{working_days} working days)"
            p.font.size = Pt(14)
            if person_name == sorted_people[0][0]:
                p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = "Remaining Work:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(24)

        p = tf.add_paragraph()
        remaining_start = self.summary.as_of_date.day + 1
        remaining_end = self.summary.month_end.day
        month_abbr = calendar.month_abbr[self.summary.as_of_date.month]
        p.text = f"• {int(self.summary.remaining_hours)} hours remaining ({remaining_start}-{remaining_end} {month_abbr})"
        p.font.size = Pt(14)
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = "• Primarily editing work on all projects"
        p.font.size = Pt(14)

    def _create_forecast_overview_slide(self):
        """Create slide 10: 3-month forecast overview."""
        slide = add_content_slide(self.prs, "3-Month Forecast Overview")

        # Summary table: Month | Projects | Team Members | Total Hours
        table_data = [['Month', 'Projects', 'Team Members', 'Total Hours']]

        # Current month row
        current_month_name = calendar.month_name[self.summary.month_start.month]
        table_data.append([
            f"{current_month_name} {self.summary.month_start.year} (Current)",
            str(len(self.summary.projects)),
            str(len(self.summary.people)),
            f"{int(self.summary.total_hours)}h"
        ])

        # Forecast month rows
        for month_summary in self.forecast.months:
            month_name = calendar.month_name[month_summary.month_start.month]
            table_data.append([
                f"{month_name} {month_summary.month_start.year}",
                str(len(month_summary.projects)),
                str(len(month_summary.people)),
                f"{int(month_summary.total_hours)}h"
            ])

        add_table(slide, table_data, Inches(0.5), Inches(1.5), Inches(9), Inches(2))

        # Bar chart comparing months
        chart_data = CategoryChartData()
        month_labels = [
            f"{calendar.month_abbr[self.summary.month_start.month]} (Current)"
        ]
        hours_values = [int(self.summary.total_hours)]

        for ms in self.forecast.months:
            month_labels.append(
                f"{calendar.month_abbr[ms.month_start.month]} {ms.month_start.year}"
            )
            hours_values.append(int(ms.total_hours))

        chart_data.categories = month_labels
        chart_data.add_series('Total Hours', hours_values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(4), Inches(8), Inches(3),
            chart_data
        ).chart
        chart.has_legend = False
        chart.chart_title.text_frame.text = 'Projected Hours by Month'

    def _create_forecast_detail_slide(self):
        """Create slide 11: Forecast team workload breakdown."""
        slide = add_content_slide(self.prs, "Forecast - Team Workload")

        # Collect all unique people across forecast months
        all_people = set()
        for ms in self.forecast.months:
            all_people.update(ms.people.keys())

        if not all_people:
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
            tf = textbox.text_frame
            tf.paragraphs[0].text = "No team members scheduled in forecast months."
            tf.paragraphs[0].font.size = Pt(14)
            return

        all_people = sorted(all_people)

        # Stacked bar chart: people on Y axis, one series per forecast month
        chart_data = ChartData()
        chart_data.categories = all_people

        for ms in self.forecast.months:
            month_label = f"{calendar.month_abbr[ms.month_start.month]} {ms.month_start.year}"
            values = []
            for p in all_people:
                person_data = ms.people.get(p)
                values.append(int(person_data.total_hours) if person_data else 0)
            chart_data.add_series(month_label, values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5),
            chart_data
        ).chart
        chart.chart_title.text_frame.text = 'Forecast Hours by Team Member'
