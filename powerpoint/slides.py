"""Reusable PowerPoint slide creation functions."""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from ..config import HEADER_COLOR, ALT_ROW_COLOR, WHITE_COLOR


def add_title_slide(prs, title, subtitle):
    """
    Add a title slide to the presentation.

    Args:
        prs: Presentation object
        title: Title text
        subtitle: Subtitle text

    Returns:
        The created slide
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle
    return slide


def add_content_slide(prs, title):
    """
    Add a content slide with a title.

    Args:
        prs: Presentation object
        title: Slide title

    Returns:
        The created slide
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    return slide


def add_table(slide, data, left, top, width, height):
    """
    Add a formatted table to a slide.

    Args:
        slide: Slide object
        data: 2D list of table data
        left: Left position
        top: Top position
        width: Table width
        height: Table height

    Returns:
        The created table
    """
    rows = len(data)
    cols = len(data[0])

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for col_idx in range(cols):
        table.columns[col_idx].width = Inches(width.inches / cols)

    # Fill table
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)

            # Format header row
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*HEADER_COLOR)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(*WHITE_COLOR)
                        run.font.bold = True
                        run.font.size = Pt(11)
            else:
                # Alternate row colours
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*ALT_ROW_COLOR)

    return table


def add_bar_chart(slide, categories, values, chart_title, left, top, width, height):
    """
    Add a bar chart to a slide.

    Args:
        slide: Slide object
        categories: List of category names
        values: List of values
        chart_title: Chart title
        left: Left position
        top: Top position
        width: Chart width
        height: Chart height

    Returns:
        The created chart
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Hours', values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    ).chart

    chart.has_legend = False
    chart.chart_title.text_frame.text = chart_title

    return chart


def add_pie_chart(slide, categories, values, chart_title, left, top, width, height):
    """
    Add a pie chart to a slide.

    Args:
        slide: Slide object
        categories: List of category names
        values: List of values
        chart_title: Chart title
        left: Left position
        top: Top position
        width: Chart width
        height: Chart height

    Returns:
        The created chart
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Hours', values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    ).chart

    chart.has_legend = True
    chart.chart_title.text_frame.text = chart_title

    return chart
